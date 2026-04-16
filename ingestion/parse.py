#!/usr/bin/env python3
"""
Eidetic ingestion parser.

Usage:
    python parse.py <doc_id> <file_path> <db_path>

Dispatches to the correct parser based on file extension, inserts
document_fragments into the SQLite database, and updates the document status.
"""

import sys
import os
import json
import sqlite3
import uuid
import traceback
from pathlib import Path
from datetime import datetime


# ---------------------------------------------------------------------------
# DB helpers
# ---------------------------------------------------------------------------

def get_db(db_path: str) -> sqlite3.Connection:
    conn = sqlite3.connect(db_path)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    conn.execute("PRAGMA journal_mode = WAL")
    return conn


def set_status(conn: sqlite3.Connection, doc_id: str, status: str, error: str | None = None):
    conn.execute(
        "UPDATE documents SET status = ?, error_message = ?, updated_at = ? WHERE id = ?",
        (status, error, int(datetime.now().timestamp() * 1000), doc_id),
    )
    conn.commit()


def insert_fragments(conn: sqlite3.Connection, doc_id: str, fragments: list[dict]):
    now = int(datetime.now().timestamp() * 1000)
    for frag in fragments:
        conn.execute(
            """INSERT INTO document_fragments
               (id, document_id, text, fragment_type,
                page_number, slide_number, slide_title,
                sheet_name, row_number,
                vendor, receipt_date, total,
                metadata_json, created_at)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (
                str(uuid.uuid4()),
                doc_id,
                frag.get("text", ""),
                frag.get("fragment_type", "text"),
                frag.get("page_number"),
                frag.get("slide_number"),
                frag.get("slide_title"),
                frag.get("sheet_name"),
                frag.get("row_number"),
                frag.get("vendor"),
                frag.get("receipt_date"),
                frag.get("total"),
                frag.get("metadata_json"),
                now,
            ),
        )
    conn.execute(
        "UPDATE documents SET fragment_count = ?, updated_at = ? WHERE id = ?",
        (len(fragments), now, doc_id),
    )
    conn.commit()


# ---------------------------------------------------------------------------
# Parsers
# ---------------------------------------------------------------------------

def parse_pdf_or_docx(file_path: str) -> list[dict]:
    """Parse PDF and Office documents via unstructured."""
    try:
        from unstructured.partition.auto import partition
    except ImportError:
        raise RuntimeError(
            "unstructured is not installed. Run: pip install unstructured[pdf,docx]"
        )

    elements = partition(filename=file_path)
    fragments = []
    for el in elements:
        text = str(el).strip()
        if not text:
            continue
        meta = el.metadata
        fragments.append(
            {
                "text": text,
                "fragment_type": type(el).__name__.lower(),
                "page_number": getattr(meta, "page_number", None),
            }
        )
    return fragments


def parse_pptx(file_path: str) -> list[dict]:
    """Parse PowerPoint files via python-pptx."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        raise RuntimeError(
            "python-pptx is not installed. Run: pip install python-pptx"
        )

    prs = Presentation(file_path)
    fragments = []

    for slide_num, slide in enumerate(prs.slides, 1):
        # Identify slide title
        slide_title = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:  # title placeholder
                slide_title = shape.text_frame.text.strip() or None
                break

        # Extract text from all shapes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            fragments.append(
                {
                    "text": text,
                    "fragment_type": "slide_text",
                    "slide_number": slide_num,
                    "slide_title": slide_title,
                }
            )

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                fragments.append(
                    {
                        "text": notes_text,
                        "fragment_type": "speaker_notes",
                        "slide_number": slide_num,
                        "slide_title": slide_title,
                    }
                )

    return fragments


def parse_xlsx(file_path: str) -> list[dict]:
    """Parse Excel spreadsheets via openpyxl."""
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError(
            "openpyxl is not installed. Run: pip install openpyxl"
        )

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    fragments = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        headers: list[str] | None = None

        for row_num, row in enumerate(ws.iter_rows(values_only=True), 1):
            if all(c is None for c in row):
                continue

            if headers is None:
                # First non-empty row becomes the header
                headers = [
                    str(c).strip() if c is not None else f"col_{i}"
                    for i, c in enumerate(row)
                ]
                continue

            row_dict = {
                headers[i]: row[i]
                for i in range(min(len(headers), len(row)))
                if row[i] is not None
            }
            if not row_dict:
                continue

            text = ", ".join(f"{k}: {v}" for k, v in row_dict.items())
            fragments.append(
                {
                    "text": text,
                    "fragment_type": "spreadsheet_row",
                    "sheet_name": sheet_name,
                    "row_number": row_num,
                    "metadata_json": json.dumps(row_dict, default=str),
                }
            )

    return fragments


def parse_image(file_path: str) -> list[dict]:
    """Parse receipt images via pytesseract + Pillow."""
    try:
        import pytesseract
        from PIL import Image, ImageFilter
    except ImportError:
        raise RuntimeError(
            "pytesseract or Pillow is not installed. "
            "Run: pip install pytesseract Pillow\n"
            "Also install Tesseract OCR: https://github.com/UB-Mannheim/tesseract/wiki"
        )

    img = Image.open(file_path).convert("L")  # greyscale
    img = img.filter(ImageFilter.SHARPEN)
    raw_text = pytesseract.image_to_string(img).strip()

    if not raw_text:
        raise RuntimeError("Tesseract returned empty text — check that Tesseract is installed and the image is readable")

    fragments: list[dict] = [{"text": raw_text, "fragment_type": "ocr_text"}]

    # Heuristic field extraction
    import re

    total_match = re.search(
        r"(?:total|amount\s*due|grand\s*total)[:\s]*\$?\s*([\d,]+\.?\d*)",
        raw_text,
        re.IGNORECASE,
    )
    subtotal_match = re.search(
        r"(?:subtotal|sub\s*total)[:\s]*\$?\s*([\d,]+\.?\d*)",
        raw_text,
        re.IGNORECASE,
    )
    tax_match = re.search(
        r"(?:tax|hst|gst|vat)[:\s]*\$?\s*([\d,]+\.?\d*)",
        raw_text,
        re.IGNORECASE,
    )
    date_match = re.search(
        r"\b(\d{1,2}[/\-\.]\d{1,2}[/\-\.]\d{2,4})\b",
        raw_text,
    )

    lines = [l.strip() for l in raw_text.splitlines() if l.strip()]
    vendor = lines[0] if lines else None

    fields: dict[str, str | None] = {
        "vendor": vendor,
        "receipt_date": date_match.group(1) if date_match else None,
        "total": total_match.group(1) if total_match else None,
        "subtotal": subtotal_match.group(1) if subtotal_match else None,
        "tax": tax_match.group(1) if tax_match else None,
    }
    # Remove None values for cleaner storage
    fields = {k: v for k, v in fields.items() if v is not None}

    if fields:
        fragments.append(
            {
                "text": json.dumps(fields),
                "fragment_type": "receipt_fields",
                "vendor": fields.get("vendor"),
                "receipt_date": fields.get("receipt_date"),
                "total": fields.get("total"),
                "metadata_json": json.dumps(fields),
            }
        )

    return fragments


# ---------------------------------------------------------------------------
# Phase 8 — change detection: mark wiki pages dirty after re-parse
# ---------------------------------------------------------------------------

def mark_dirty_pages(conn: sqlite3.Connection, doc_id: str):
    """Find wiki pages whose source_doc_ids include this document and mark them dirty."""
    now = int(datetime.now().timestamp() * 1000)
    pages = conn.execute(
        "SELECT id, slug, source_doc_ids FROM wiki_pages"
    ).fetchall()

    dirty_count = 0
    for page in pages:
        source_ids = json.loads(page["source_doc_ids"] or "[]")
        if doc_id in source_ids:
            conn.execute(
                "UPDATE wiki_pages SET dirty = 1, updated_at = ? WHERE id = ?",
                (now, page["id"]),
            )
            # Enqueue a rebuild job (skip if one is already pending)
            existing_job = conn.execute(
                "SELECT id FROM job_queue WHERE target_id = ? AND job_type = 'rewiki' AND status = 'pending'",
                (page["slug"],),
            ).fetchone()
            if not existing_job:
                conn.execute(
                    """INSERT INTO job_queue (id, job_type, target_id, status, reason, created_at, updated_at)
                       VALUES (?, 'rewiki', ?, 'pending', ?, ?, ?)""",
                    (str(uuid.uuid4()), page["slug"], f"document {doc_id} re-parsed", now, now),
                )
            dirty_count += 1

    if dirty_count:
        conn.commit()
        print(f"[parse] Marked {dirty_count} wiki page(s) as dirty")


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

PARSERS = {
    "pdf": parse_pdf_or_docx,
    "docx": parse_pdf_or_docx,
    "pptx": parse_pptx,
    "xlsx": parse_xlsx,
    "image": parse_image,
}


def main():
    if len(sys.argv) != 4:
        print(f"Usage: {sys.argv[0]} <doc_id> <file_path> <db_path>", file=sys.stderr)
        sys.exit(1)

    doc_id = sys.argv[1]
    file_path = sys.argv[2]
    db_path = sys.argv[3]

    conn = get_db(db_path)

    try:
        # Determine file type from DB record
        row = conn.execute(
            "SELECT file_type FROM documents WHERE id = ?", (doc_id,)
        ).fetchone()
        if row is None:
            print(f"Document {doc_id} not found in DB", file=sys.stderr)
            sys.exit(1)

        file_type = row["file_type"]
        parser = PARSERS.get(file_type)
        if parser is None:
            raise RuntimeError(f"No parser registered for file type: {file_type}")

        if not os.path.isfile(file_path):
            raise RuntimeError(f"File not found: {file_path}")

        print(f"[parse] Parsing {file_path} as {file_type}…")
        fragments = parser(file_path)
        print(f"[parse] Extracted {len(fragments)} fragment(s)")

        insert_fragments(conn, doc_id, fragments)
        set_status(conn, doc_id, "processed")

        # Phase 8 — mark wiki pages referencing this document as dirty
        mark_dirty_pages(conn, doc_id)

        print(f"[parse] Done — document {doc_id} status: processed")

    except Exception as exc:
        error_msg = str(exc)
        print(f"[parse] ERROR: {error_msg}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        set_status(conn, doc_id, "failed", error_msg[:2000])
        sys.exit(1)
    finally:
        conn.close()


if __name__ == "__main__":
    main()
